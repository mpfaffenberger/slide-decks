#!/usr/bin/env python3
"""Build the RVA Data+AI Summit presentation as .pptx (opens in Keynote)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Dark theme colors
BG = RGBColor(0x11, 0x11, 0x1F)  # darker charcoal-navy — projector safe
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x2D, 0x9C, 0xDB)  # cerulean blue — trust, authority, calm
MUTED = RGBColor(0xA0, 0xA8, 0xBC)  # soft gray
WARN = RGBColor(0xFF, 0xD7, 0x66)    # warm yellow

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def set_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_text(slide, left, top, width, height, text, size=28, color=WHITE,
             bold=False, alignment=PP_ALIGN.LEFT, font_name="Menlo"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullets(slide, left, top, width, height, items, size=22, color=WHITE):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Menlo"
        p.space_after = Pt(10)
    return txBox


# ──────────────────────────────────────────
# SLIDE 1: Title
# ──────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(s)
add_text(s, 1, 1.5, 11, 1.5,
         "Open-Source AI and the\nNew Data Control Plane",
         size=48, bold=True, color=WHITE, font_name="Avenir Next")
add_text(s, 1, 3.8, 11, 0.8,
         "RVA Data+AI Summit  ·  March 26, 2026",
         size=24, color=MUTED)
add_text(s, 1, 4.6, 11, 0.8,
         "Mike Pfaffenberger  &  Mari",
         size=24, color=ACCENT)

# ──────────────────────────────────────────
# SLIDE 2: The Hook
# ──────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 1, 1, 11, 2,
         '"Your AI agent just dropped a 200-line\nSQL refactor. Do you trust it?"',
         size=44, bold=True, color=WARN, font_name="Avenir Next")
add_text(s, 1, 4, 11, 2,
         "Quick poll:\n"
         "✋  How many of you have an AI agent in production?\n"
         "✋  How many have a review process for AI-generated analyses?\n\n"
         "The gap between those hands is what this talk is about.",
         size=24, color=WHITE)

# ──────────────────────────────────────────
# SLIDE 3: The Trust Problem
# ──────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 1, 0.8, 11, 1, "The Trust Problem", size=42, bold=True, color=ACCENT, font_name="Avenir Next")
add_text(s, 1, 2, 11, 1,
         "Agentic AI has moved from suggesting to executing.\nThis breaks every existing trust model we have.",
         size=26, color=WHITE)

# ──────────────────────────────────────────
# SLIDE 4: Failure Mode 1
# ──────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 1, 0.8, 11, 1, "The Phantom Dashboard", size=40, bold=True, color=WARN, font_name="Avenir Next")
add_text(s, 1, 2.2, 11, 4,
         "Agent builds a beautiful analysis.\n\n"
         "Nobody knows where the data came from.\n\n"
         "It gets screenshot'd into Slack and becomes gospel.\n\n"
         "No lineage. No review.",
         size=26, color=WHITE)
add_text(s, 1, 5.8, 6, 0.6, "Missing layer: Governance & Review", size=20, color=MUTED)

# ──────────────────────────────────────────
# SLIDE 5: Failure Mode 2
# ──────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 1, 0.8, 11, 1, "The Black Box Refactor", size=40, bold=True, color=WARN, font_name="Avenir Next")
add_text(s, 1, 2.2, 11, 4,
         "Agent rewrites a pipeline. It works.\n\n"
         "Six months later it breaks.\n\n"
         "Nobody knows why those decisions were made.",
         size=26, color=WHITE)
add_text(s, 1, 5.8, 8, 0.6, "Missing layer: Observability + Transparent Execution", size=20, color=MUTED)

# ──────────────────────────────────────────
# SLIDE 6: Failure Mode 3
# ──────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 1, 0.8, 11, 1, "The Hallucinated Insight", size=40, bold=True, color=WARN, font_name="Avenir Next")
add_text(s, 1, 2.2, 11, 4,
         "Agent produces a compelling but subtly wrong analysis.\n\n"
         "It gets reused across 3 teams\nbefore someone catches it.",
         size=26, color=WHITE)
add_text(s, 1, 5.8, 8, 0.6, "Missing layer: Versioning, Review, Attribution", size=20, color=MUTED)

# ──────────────────────────────────────────
# SLIDE 7: The Trust Stack
# ──────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 1, 0.5, 11, 1, "The Trust Stack", size=42, bold=True, color=ACCENT, font_name="Avenir Next")

layers = [
    ("Trusted Insight", "Versioned, reviewed, reusable"),
    ("Governance & Review", "Who approved it? What changed?"),
    ("Observability", "Can I see what the agent did?"),
    ("Transparent Execution", "Open-source, inspectable"),
    ("AI Model", "The LLM (commodity layer)"),
]

for i, (label, desc) in enumerate(layers):
    y = 1.8 + i * 1.0
    # box
    shape = s.shapes.add_shape(
        1, Inches(2), Inches(y), Inches(9), Inches(0.85)
    )
    shape.fill.solid()
    # gradient from accent to darker
    alpha = 255 - (i * 40)
    shape.fill.fore_color.rgb = RGBColor(
        max(0x12, 0x2D - i * 8),
        max(0x3A, 0x9C - i * 20),
        max(0x6A, 0xDB - i * 25),
    )
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"{label}  —  {desc}"
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.font.name = "Helvetica Neue"
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

add_text(s, 1, 6.8, 11, 0.5,
         "Most vendors sell you the bottom layer. Nobody's building the top.",
         size=20, color=MUTED, alignment=PP_ALIGN.CENTER)

# ──────────────────────────────────────────
# SLIDE 8: Transparency — Code Puppy
# ──────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 1, 0.8, 11, 1, "Transparency as Foundation", size=42, bold=True, color=ACCENT, font_name="Avenir Next")
add_text(s, 1, 2, 11, 1,
         '"If you can\'t inspect it, you can\'t trust it."',
         size=30, bold=True, color=WARN)
add_bullets(s, 1, 3.5, 11, 3, [
    "Open source isn't ideology — it's falsifiability",
    "Closed agents are unfalsifiable: you can't verify what you can't see",
    "Every action has a hook for observation and intervention",
    "pre_tool_call → post_tool_call → file_permission → run_shell_command",
], size=24)

# ──────────────────────────────────────────
# SLIDE 9: Live Demo
# ──────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 2, 2.5, 9, 2,
         "🐶  Live Demo",
         size=54, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER, font_name="Avenir Next")
add_text(s, 2, 4.2, 9, 1,
         "An open-source AI agent executing real data work —\nwith every decision observable and interceptable.",
         size=24, color=MUTED, alignment=PP_ALIGN.CENTER)

# ──────────────────────────────────────────
# SLIDE 10: Demo Punchline / Pivot
# ──────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 1, 1.5, 11, 3,
         "You just watched an AI agent do real work.\n"
         "You could see every decision it made.\n\n"
         "But where does this output go?\n"
         "Who reviews it?\n"
         "Who owns it six months from now?",
         size=32, color=WHITE)
add_text(s, 1, 5.5, 11, 1,
         "Transparency alone can't answer these questions.",
         size=28, bold=True, color=WARN)

# ──────────────────────────────────────────
# SLIDE 11: Transparency ≠ Trust
# ──────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 1, 0.8, 11, 1, "Transparency ≠ Trust", size=42, bold=True, color=ACCENT, font_name="Avenir Next")
add_text(s, 1, 2.2, 11, 1,
         "We solved this for code 20 years ago — git + PR reviews.\n"
         "We haven't solved it for data work.",
         size=26, color=WHITE)
add_bullets(s, 1, 3.8, 11, 3, [
    "Agents produce artifacts. Teams need insights.",
    "An artifact is a file. An insight is versioned, attributed, reviewed knowledge.",
    "AI-generated analyses disappear into notebooks, dashboards, Slack threads.",
    "This problem exists with or without AI. AI just made it urgent.",
], size=24)

# ──────────────────────────────────────────
# SLIDE 12: The Control Plane
# ──────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 1, 0.8, 11, 1, "The Missing Layer: A Control Plane for Data Work", size=38, bold=True, color=ACCENT, font_name="Avenir Next")
add_bullets(s, 1, 2.2, 11, 4.5, [
    "Capture — AI and human analyses land in one place",
    "Version — Every change is tracked, diffable, attributable",
    "Review — Approval workflows before insights become \"official\"",
    "Reuse — Insights become organizational knowledge",
    "Govern — Lineage, access control, audit trails",
], size=26)

# ──────────────────────────────────────────
# SLIDE 13: Shared Design Principles
# ──────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 1, 0.8, 11, 1, "Two Angles, One Belief", size=42, bold=True, color=ACCENT, font_name="Avenir Next")
add_text(s, 1, 2.2, 11, 3,
         "Code Puppy tackles the execution layer —\n"
         "making agent behavior transparent and inspectable.\n\n"
         "Ara tackles the knowledge layer —\n"
         "making sure insights are captured, reviewed, and reusable.\n\n"
         "Independent projects. Shared core belief:",
         size=26, color=WHITE)
add_text(s, 1, 5.5, 11, 1,
         "Trust has to be designed in, not bolted on.",
         size=32, bold=True, color=WARN, alignment=PP_ALIGN.CENTER)

# ──────────────────────────────────────────
# SLIDE 14: T.R.U.S.T. Framework
# ──────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 1, 0.5, 11, 1, "T.R.U.S.T.", size=48, bold=True, color=ACCENT, font_name="Avenir Next")
add_text(s, 1, 1.3, 11, 0.5,
         "A checklist for evaluating any tool in your data stack.",
         size=22, color=MUTED)

items = [
    ("T", "Transparency", "Can I see what happened?"),
    ("R", "Review", "Did a human validate it?"),
    ("U", "Understanding", "Is the context preserved?"),
    ("S", "Stability", "Can I reproduce it?"),
    ("T", "Traceability", "Can I audit the full chain?"),
]

for i, (letter, principle, question) in enumerate(items):
    y = 2.2 + i * 0.95
    add_text(s, 1.5, y, 0.8, 0.8, letter, size=36, bold=True, color=ACCENT)
    add_text(s, 2.5, y, 4, 0.8, principle, size=28, bold=True, color=WHITE)
    add_text(s, 6.5, y + 0.05, 5.5, 0.8, question, size=24, color=MUTED)

# ──────────────────────────────────────────
# SLIDE 15: Practical Advice
# ──────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 1, 0.8, 11, 1, "What You Can Do Monday Morning", size=42, bold=True, color=ACCENT, font_name="Avenir Next")
add_bullets(s, 1, 2.2, 11, 4.5, [
    "Start with transparency — use open-source agents you can inspect",
    "Add review workflows for AI-generated outputs — don't auto-ship",
    "Invest in a durable home for insights — not Slack threads and Google Docs",
    "Treat AI outputs like PRs: valuable, but not trusted until reviewed",
    "Evaluate every tool against the T.R.U.S.T. checklist",
], size=26)

# ──────────────────────────────────────────
# SLIDE 16: Closing
# ──────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 1, 1.5, 11, 2,
         '"We don\'t have an AI trust problem.\n'
         'We have a data governance problem\n'
         'that AI made urgent."',
         font_name="Avenir Next",
         size=40, bold=True, color=WARN, alignment=PP_ALIGN.CENTER)
add_text(s, 1, 4.5, 11, 1,
         "Trust has to be designed in, not bolted on.",
         size=28, color=ACCENT, alignment=PP_ALIGN.CENTER)

# ──────────────────────────────────────────
# SLIDE 17: Thank You
# ──────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
add_text(s, 2, 2, 9, 1.5,
         "Thank You",
         size=54, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER, font_name="Avenir Next")
add_text(s, 2, 3.8, 9, 1,
         "Mike Pfaffenberger  &  Mari",
         size=28, color=ACCENT, alignment=PP_ALIGN.CENTER)
add_text(s, 2, 5, 9, 1,
         "Questions?",
         size=32, color=MUTED, alignment=PP_ALIGN.CENTER)

# Save
out = "RVA_Data_AI_Summit.pptx"
prs.save(out)
print(f"✅ Saved {out} — open in Keynote!")
